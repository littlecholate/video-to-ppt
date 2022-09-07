import numpy as np
import tkinter as tk
import re
import os
import cv2
import threading
from pytube import YouTube
from pptx import Presentation
from pptx.util import Inches
from tkinter.scrolledtext import ScrolledText

def onProgress(stream, chunk, remains):
    total = stream.filesize # 取得完整尺寸
    percent = (total-remains)/total*100 # 減去剩餘尺寸(剩餘尺寸會抓取存取的檔案大小)
    print(f'下載中 {percent:05.2f} %', end='\r') # 顯示進度，\r 表示不換行，在同一行更新

def MSE(img1, img2):
        squared_diff = (img1-img2)**2
        summed = np.sum(squared_diff)
        num_pix = img1.shape[0]*img1.shape[1]
        err = summed/num_pix
        return err

def thread_it(func, *args):
    t = threading.Thread(target=func, args=args)
    t.setDaemon(True)
    t.start()

def run():
    stext.insert(tk.END, '開始下載影片...\n')
    yt = YouTube(input_url.get())
    yt.streams.filter().get_by_resolution('720p').download(filename='input.mp4')
    stext.insert(tk.END, '下載成功\n\n')

    stext.insert(tk.END, '正在逐幀轉成圖片，會要一點時間...\n')
    if not os.path.isdir('output'):
        os.mkdir('output')
    os.system('ffmpeg -loglevel error -i input.mp4 -vf fps=1 output/out%d.png')
    stext.insert(tk.END, '轉成圖片完成\n\n')

    stext.insert(tk.END, '正在刪除重複的圖片...\n')
    fileNum = len([x for x in os.listdir('output/')]) # read images
    imgPrev = cv2.imread('output/out1.png') # first image
    for i in range(fileNum-1):
        filePath = 'output/out'+ str(i+2) +'.png'
        img = cv2.imread(filePath)
        grayA = cv2.cvtColor(imgPrev, cv2.COLOR_BGR2GRAY)
        grayB = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)
        # check threshold and delete images
        if MSE(grayA,grayB) < 4:
            os.remove(filePath)
        #print(str(i) + ' MSE: ' + str(MSE(grayA,grayB)))
        imgPrev = img
    stext.insert(tk.END, '刪除重複圖片完成\n\n')

    stext.insert(tk.END, '正在製作投影片...\n')
    prs = Presentation() 
    prs.slide_height = Inches(9)
    prs.slide_width = Inches(16)
    blank_slide_layout = prs.slide_layouts[0]
    # sort the images by name
    list_dir = os.listdir('output/')
    list_dir = sorted(list_dir, key=lambda s: int(re.search(r'\d+', s).group()))
    for f in list_dir:
        slide = prs.slides.add_slide(blank_slide_layout)
        pic = slide.shapes.add_picture('output/' + f, 0, 0, Inches(16), Inches(9))
        os.remove('output/' + f)
    prs.save("result.pptx")
    os.rmdir('output')
    stext.insert(tk.END, '投影片製作完成\n')

# make the main GUI
win = tk.Tk()
win.title("Youtube Video to Powerpoint")

# Detect screen size, set window's size and move the window to the center.
width = 400
height = 400
screenwidth = win.winfo_screenwidth()  
screenheight = win.winfo_screenheight() 
alignstr = '%dx%d+%d+%d' % (width, height, (screenwidth-width)/2, (screenheight-height)/2)   
win.geometry(alignstr)
win.resizable(0,0)

instruction = tk.Label(text='輸入影片網址 : ', font=('楷體', 12))
instruction.pack(padx=10, pady=10, anchor=tk.W)

input_url = tk.Entry()
input_url.pack(ipadx=120, ipady=6, padx=10, anchor=tk.W)

button = tk.Button(text='開始執行', font=('楷體',12), command=lambda :thread_it(run))
button.pack(pady=12)

stext = ScrolledText(width=60, height=18, background='#F7F3EC')
stext.pack(padx=20)

win.mainloop()



